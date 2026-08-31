#!/usr/bin/env python3
"""Generate deterministic, non-sensitive multimodal acceptance fixtures."""

from __future__ import annotations

import csv
import json
import math
import shutil
import struct
import subprocess
import wave
import zipfile
from email.message import EmailMessage
from pathlib import Path


FACT = "NexusOne is an enterprise knowledge platform released in 2026."


def run(command: list[str]) -> None:
    subprocess.run(command, check=True, capture_output=True, timeout=180)


def write_text_pdf(path: Path) -> None:
    """Write a tiny standards-compliant two-page PDF with embedded text."""
    streams = [
        b"BT /F1 18 Tf 72 760 Td (NexusOne Product Manual) Tj 0 -30 Td (Enterprise knowledge platform released in 2026.) Tj ET",
        b"BT /F1 14 Tf 72 760 Td (Data sources include REST API, Git, databases and email.) Tj ET",
    ]
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R 4 0 R] /Count 2 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 6 0 R >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 7 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(streams[0]), streams[0]),
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(streams[1]), streams[1]),
    ]
    payload = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, value in enumerate(objects, 1):
        offsets.append(len(payload))
        payload.extend(f"{number} 0 obj\n".encode())
        payload.extend(value)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode())
    payload.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    path.write_bytes(payload)


def generate(root: Path) -> dict[str, Path]:
    from docx import Document
    from ebooklib import epub
    from openpyxl import Workbook
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation

    if root.exists():
        shutil.rmtree(root)
    root.mkdir(parents=True)
    files: dict[str, Path] = {}

    def add(name: str, text: str) -> Path:
        path = root / name
        path.write_text(text, encoding="utf-8")
        files[name] = path
        return path

    add("fact.txt", FACT)
    add("fact.md", f"# Product\n\n{FACT}\n")
    add("fact.py", f'PRODUCT_FACT = "{FACT}"\n')
    add("fact.json", json.dumps({"product": "NexusOne", "year": 2026}, ensure_ascii=False))
    add("fact.jsonl", '{"product":"NexusOne"}\n{"year":2026}\n')
    add("fact.xml", "<product><name>NexusOne</name><year>2026</year></product>")
    add("fact.html", f"<!doctype html><title>NexusOne</title><p>{FACT}</p>")
    with (root / "fact.csv").open("w", encoding="utf-8", newline="") as stream:
        writer = csv.writer(stream)
        writer.writerow(["product", "year", "positioning"])
        writer.writerow(["NexusOne", 2026, "enterprise knowledge platform"])
    files["fact.csv"] = root / "fact.csv"

    document = Document()
    document.add_heading("NexusOne Product Manual", 0)
    document.add_paragraph(FACT)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text, table.cell(0, 1).text = "Capability", "Status"
    table.cell(1, 0).text, table.cell(1, 1).text = "Hybrid retrieval", "Supported"
    document.save(root / "fact.docx")
    files["fact.docx"] = root / "fact.docx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "NexusOne"
    slide.placeholders[1].text = FACT
    presentation.save(root / "fact.pptx")
    files["fact.pptx"] = root / "fact.pptx"

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Capabilities"
    sheet.append(["Product", "Year", "Positioning"])
    sheet.append(["NexusOne", 2026, "enterprise knowledge platform"])
    workbook.save(root / "fact.xlsx")
    files["fact.xlsx"] = root / "fact.xlsx"

    if shutil.which("libreoffice"):
        legacy_targets = {
            "fact.docx": ("doc:MS Word 97", "fact.doc"),
            "fact.pptx": ("ppt:MS PowerPoint 97", "fact.ppt"),
            "fact.xlsx": ("xls:MS Excel 97", "fact.xls"),
        }
        for source, (target_format, target_name) in legacy_targets.items():
            run([
                "libreoffice", "--headless", "--nologo", "--nodefault", "--nofirststartwizard",
                "--convert-to", target_format, "--outdir", str(root), str(root / source),
            ])
            if (root / target_name).exists():
                files[target_name] = root / target_name

    write_text_pdf(root / "fact.pdf")
    files["fact.pdf"] = root / "fact.pdf"

    font_path = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
    font = ImageFont.truetype(str(font_path), 34) if font_path.exists() else ImageFont.load_default()
    image = Image.new("RGB", (1400, 500), "white")
    drawing = ImageDraw.Draw(image)
    drawing.text((60, 100), "NexusOne OCR Fact", fill="black", font=font)
    drawing.text((60, 180), "Enterprise Knowledge Platform 2026", fill="black", font=font)
    image.save(root / "fact.png")
    image.save(root / "fact.jpg", quality=95)
    image.save(root / "scanned.pdf", "PDF", resolution=150)
    files.update({name: root / name for name in ("fact.png", "fact.jpg", "scanned.pdf")})

    email = EmailMessage()
    email["Subject"] = "NexusOne release"
    email["From"] = "knowledge@example.test"
    email["To"] = "group@example.test"
    email.set_content(FACT)
    email.add_attachment(
        b"NexusOne supports hybrid retrieval.",
        maintype="text",
        subtype="plain",
        filename="capabilities.txt",
    )
    (root / "fact.eml").write_bytes(email.as_bytes())
    files["fact.eml"] = root / "fact.eml"

    book = epub.EpubBook()
    book.set_identifier("nexusone-fixture")
    book.set_title("NexusOne Fixture")
    book.set_language("en")
    chapter = epub.EpubHtml(title="Positioning", file_name="positioning.xhtml", lang="en")
    chapter.content = f"<h1>NexusOne</h1><p>{FACT}</p>"
    book.add_item(chapter)
    book.toc = (epub.Link(chapter.file_name, chapter.title, "positioning"),)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(root / "fact.epub"), book)
    files["fact.epub"] = root / "fact.epub"

    sample_rate, duration = 16_000, 1.5
    with wave.open(str(root / "fact.wav"), "wb") as wav:
        wav.setparams((1, 2, sample_rate, int(sample_rate * duration), "NONE", "not compressed"))
        frames = bytearray()
        for index in range(int(sample_rate * duration)):
            value = int(0.25 * 32767 * math.sin(2 * math.pi * 440 * index / sample_rate))
            frames.extend(struct.pack("<h", value))
        wav.writeframes(bytes(frames))
    run(["ffmpeg", "-y", "-loglevel", "error", "-i", str(root / "fact.wav"), str(root / "fact.mp3")])
    run([
        "ffmpeg", "-y", "-loglevel", "error", "-f", "lavfi", "-i", "color=c=blue:s=320x240:d=2",
        "-f", "lavfi", "-i", "sine=frequency=440:duration=2", "-shortest", "-pix_fmt", "yuv420p",
        str(root / "fact.mp4"),
    ])
    files.update({name: root / name for name in ("fact.wav", "fact.mp3", "fact.mp4")})

    try:
        import pyarrow as pa
        import pyarrow.feather as feather
        import pyarrow.parquet as parquet

        columnar = pa.table({"product": ["NexusOne"], "year": [2026]})
        parquet.write_table(columnar, root / "fact.parquet")
        feather.write_feather(columnar, root / "fact.arrow")
        files.update({name: root / name for name in ("fact.parquet", "fact.arrow")})
    except ImportError:
        pass

    with zipfile.ZipFile(root / "fact.zip", "w", zipfile.ZIP_DEFLATED) as bundle:
        bundle.write(root / "fact.txt", "knowledge/fact.txt")
        bundle.write(root / "fact.json", "knowledge/fact.json")
    files["fact.zip"] = root / "fact.zip"
    return files


if __name__ == "__main__":
    generated = generate(Path(__file__).with_name("generated"))
    print(json.dumps({"count": len(generated), "files": sorted(generated)}, ensure_ascii=False))
