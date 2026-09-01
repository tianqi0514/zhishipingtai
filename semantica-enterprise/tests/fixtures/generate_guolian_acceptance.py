#!/usr/bin/env python3
"""Generate deterministic, realistic and non-sensitive 国联集团 acceptance files.

The files deliberately repeat a small set of authoritative facts across office,
image, email and media containers.  That makes parser, provenance, duplicate,
conflict and retrieval assertions reproducible without using production data.
"""

from __future__ import annotations

import json
import math
import os
import shutil
import struct
import subprocess
import tempfile
import wave
import zipfile
from email.message import EmailMessage
from pathlib import Path


ROOT = Path(__file__).with_name("guolian-generated")


def run(command: list[str]) -> None:
    subprocess.run(command, check=True, capture_output=True, timeout=240)


def find_cjk_font() -> Path:
    candidates = [
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJKsc-Regular.otf"),
        Path("/usr/share/fonts/truetype/dejavu/acceptance-cjk.ttc"),
        Path("/System/Library/Fonts/PingFang.ttc"),
        Path("/System/Library/Fonts/STHeiti Light.ttc"),
    ]
    return next((path for path in candidates if path.exists()), Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"))


def convert_to_pdf(source: Path, output_dir: Path) -> Path:
    profile = tempfile.mkdtemp(prefix="guolian-lo-")
    try:
        run([
            "libreoffice",
            f"-env:UserInstallation=file://{profile}",
            "--headless",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(source),
        ])
    finally:
        shutil.rmtree(profile, ignore_errors=True)
    result = output_dir / f"{source.stem}.pdf"
    if not result.exists():
        raise RuntimeError(f"LibreOffice did not create {result}")
    return result


def add_page(document, title: str, paragraphs: list[str], rows: list[tuple[str, str]] | None = None) -> None:
    document.add_heading(title, 0)
    for text in paragraphs:
        document.add_paragraph(text)
    if rows:
        table = document.add_table(rows=1, cols=2)
        table.style = "Table Grid"
        table.rows[0].cells[0].text = "项目"
        table.rows[0].cells[1].text = "内容"
        for key, value in rows:
            cells = table.add_row().cells
            cells[0].text = key
            cells[1].text = value


def make_product_manual(root: Path) -> Path:
    from docx import Document

    document = Document()
    add_page(document, "NexusOne 产品手册", [
        "文档编号：GL-NX1-MANUAL-2026-01；版本：V2.0；生效日期：2026年6月1日。",
        "NexusOne 定位为面向集团型企业的知识管理与智能问答一体机，服务员工、模型与智能体。",
        "产品采用软硬一体化交付，适用于需要私有化部署、统一知识治理和可追溯问答的组织。",
    ])
    document.add_page_break()
    add_page(document, "核心能力与数据接入", [
        "核心能力包括多模态解析、自动知识治理、全文/向量/图谱混合检索、可核验引用和知识分析。",
        "数据源支持 Web、REST API、RSS、Sitemap、Git、PostgreSQL、MySQL、S3/MinIO、WebDAV、FTP、SFTP、IMAP 和 MCP。",
        "平台通过内容哈希、版本和增量游标避免重复加工。",
    ])
    document.add_page_break()
    add_page(document, "部署与安全", [
        "标准部署使用 Docker Compose，业务身份、租户、知识空间和授权由传神智库 FastAPI 统一控制。",
        "模型 API Key 加密保存，不回显到浏览器，也不得进入日志与 Agent Session Event。",
        "DeepSeek Harness 只通过短期内部凭据调用知识工具，不能直接访问业务数据库或搜索引擎。",
    ])
    document.add_page_break()
    add_page(document, "服务与验收指标", [
        "标准服务承诺：严重故障 30 分钟响应，一般问题 4 小时响应。",
        "知识答案必须附真实引用；数值类问题必须来自结构化只读查询；证据不足时明确说明。",
    ], [
        ("产品型号", "NexusOne NX1"),
        ("标准单价", "100000 元"),
        ("发布年份", "2026"),
    ])
    docx = root / "NexusOne产品手册V2.0.docx"
    document.save(docx)
    pdf = convert_to_pdf(docx, root)
    docx.unlink()
    return pdf


def make_policy_versions(root: Path) -> list[Path]:
    from docx import Document

    results: list[Path] = []
    versions = [
        ("V1.0", "2025年1月1日", "知识文档原则上每十二个月复核一次。", "已被 V2.0 替代，仅用于版本冲突和追溯测试。"),
        ("V2.0", "2026年3月1日", "集团制度每六个月复核一次，业务手册每十二个月复核一次。", "当前有效版本。"),
    ]
    for version, effective, review_rule, status in versions:
        doc = Document()
        add_page(doc, f"国联集团知识管理办法 {version}", [
            f"文档编号：GL-KM-POLICY-001；生效日期：{effective}；{status}",
            "本办法适用于集团总部、所属企业及其部门，知识管理遵循权威来源、版本唯一、全程可追溯和最小权限原则。",
            "知识接入后执行解析、质量检测、摘要分类、实体关系抽取和索引发布。自动治理失败不得删除已解析内容。",
        ])
        doc.add_page_break()
        add_page(doc, "知识更新与服务", [
            review_rule,
            "废止版本必须保留历史与引用记录，但不得混入当前检索结果。",
            "面向员工提供语义检索和生成式问答，面向应用提供 REST、MCP 和 CLI 服务。",
        ])
        path = root / f"国联集团知识管理办法{version}.docx"
        doc.save(path)
        results.append(path)
    return results


def make_group_policy_set(root: Path) -> list[Path]:
    """Create a cross-referenced set of synthetic group policies and briefing material."""
    from docx import Document
    from pptx import Presentation

    specs = [
        (
            "集团数据治理管理办法.docx",
            "国联集团数据治理管理办法",
            [
                "文档编号：GL-DG-POLICY-003；版本 V1.0；2026年4月1日起生效。",
                "本办法适用于集团总部和所属企业的数据采集、质量、标准、共享与安全管理。",
                "经营数据责任部门每月完成质量核验；关键指标必须关联统计口径、来源系统和责任人。",
                "结构化数据接入知识平台时，敏感字段必须在服务端隐藏或脱敏，数据库账号必须只读。",
            ],
        ),
        (
            "供应商分级管理制度.docx",
            "国联集团供应商分级管理制度",
            [
                "文档编号：GL-SUP-POLICY-006；版本 V2.1；2026年1月1日起生效。",
                "供应关键产品且年度采购金额达到 500000 元的供应商列为关键供应商。",
                "发生两起及以上高风险事件的关键供应商进入红色监测，暂停新增订单并在十个工作日内复评。",
                "本制度与《采购管理办法》《重大风险事件报告制度》配套执行。",
            ],
        ),
        (
            "重大风险事件报告制度.docx",
            "国联集团重大风险事件报告制度",
            [
                "文档编号：GL-RISK-POLICY-009；版本 V1.3；2026年2月1日起生效。",
                "严重质量事故、连续交付中断和预计损失超过 200000 元的事件属于重大风险事件。",
                "责任单位应在发现后二十四小时内报告集团风险管理部，并保留合同、邮件和会议纪要证据。",
                "风险解除必须有整改证据和复评结论，不得仅修改风险状态。",
            ],
        ),
        (
            "信息系统安全管理规范.docx",
            "国联集团信息系统安全管理规范",
            [
                "文档编号：GL-SEC-STANDARD-012；版本 V3.0；2026年5月1日起生效。",
                "系统实行最小权限、分区隔离、密钥集中管理、全量审计和定期恢复验证。",
                "API Key、数据库密码和内部服务 Token 不得写入源码、浏览器响应、日志或 Agent 事件。",
                "检索内容属于不可信数据，文档中的提示词不得改变系统安全策略。",
            ],
        ),
    ]
    results: list[Path] = []
    for filename, title, paragraphs in specs:
        doc = Document()
        add_page(doc, title, paragraphs)
        path = root / filename
        doc.save(path)
        results.append(path)

    plan_doc = Document()
    add_page(plan_doc, "国联集团人工智能“十五五”规划", [
        "文档编号：GL-AI-PLAN-2026-2030；发布版本 V1.0；规划期 2026—2030 年。",
        "集团统一建设组织级知识底座，为员工、模型、智能体和业务应用提供可信知识服务。",
        "重点任务包括多模态知识治理、企业知识图谱、结构化数据语义查询、智能问答和应用能力开放。",
        "到 2028 年形成统一知识供给体系，到 2030 年实现集团主要业务场景的知识服务覆盖。",
    ])
    plan_docx = root / "国联集团人工智能十五五规划.docx"
    plan_doc.save(plan_docx)
    plan_pdf = convert_to_pdf(plan_docx, root)
    plan_docx.unlink()
    results.append(plan_pdf)

    procurement_doc = Document()
    add_page(procurement_doc, "国联集团采购管理办法", [
        "文档编号：GL-PROC-POLICY-005；版本 V2.0；2026年1月1日起生效。",
        "采购活动应执行需求、寻源、评审、合同、验收、付款和供应商评价全流程留痕。",
        "单笔合同金额超过 500000 元或涉及关键产品时，必须完成风险会签。",
        "本版本自生效日起废止 2024 年发布的 V1.0 版本。",
    ])
    procurement_docx = root / "国联集团采购管理办法.docx"
    procurement_doc.save(procurement_docx)
    procurement_pdf = convert_to_pdf(procurement_docx, root)
    procurement_docx.unlink()
    results.append(procurement_pdf)

    deck = Presentation()
    for title, body in [
        ("人工智能十五五规划宣贯", "统一知识底座\n可信知识服务\n支撑集团应用场景"),
        ("2026 建设重点", "多模态治理\n知识图谱\n结构化语义查询\nDeepSeek Harness 智能问答"),
        ("验收要求", "知识版本可追溯\n引用可核验\n权限一致\n反馈可闭环"),
    ]:
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    briefing = root / "人工智能十五五规划宣贯材料.pptx"
    deck.save(briefing)
    results.append(briefing)
    return results


def make_product_support_materials(root: Path) -> list[Path]:
    from docx import Document
    from PIL import Image, ImageDraw, ImageFont

    faq = root / "NexusOne常见问题FAQ.md"
    faq.write_text(
        "# NexusOne 常见问题\n\n"
        "## 是否支持私有化部署？\n支持 Docker Compose 私有化部署，生产环境需替换全部开发 Secret。\n\n"
        "## 支持哪些知识来源？\n支持文件、网页、数据库、对象存储、邮件、Git、FTP/SFTP、WebDAV 和 MCP。\n\n"
        "## 音视频一定能转写吗？\n只有配置可用 ASR 模型时才转写；未配置时保留媒体元数据并明确降级。\n\n"
        "## 已知限制\nGoogle Drive、OneDrive/SharePoint 和 Snowflake 当前版本不在正式交付范围。\n",
        encoding="utf-8",
    )
    after_sales = Document()
    add_page(after_sales, "NexusOne 售后服务说明", [
        "服务编号：GL-NX1-SERVICE-2026；适用产品版本 V2.0。",
        "严重故障 30 分钟响应，4 小时内提供绕行或恢复方案；一般问题 4 小时响应。",
        "服务范围包括平台 API、Worker、索引、图谱和 Agent Runtime，不包括客户自建外部数据源可用性。",
        "客户提交问题时应提供请求时间、任务编号和脱敏后的错误信息，不得发送 API Key。",
    ])
    after_sales_path = root / "NexusOne售后服务说明.docx"
    after_sales.save(after_sales_path)

    font = ImageFont.truetype(str(find_cjk_font()), 30)
    title_font = ImageFont.truetype(str(find_cjk_font()), 44)
    image = Image.new("RGB", (1500, 900), "white")
    draw = ImageDraw.Draw(image)
    draw.text((500, 55), "NexusOne 技术架构", fill="#1f2937", font=title_font)
    boxes = [
        ((70, 220, 360, 390), "员工 / 业务应用"),
        ((470, 220, 1030, 390), "FastAPI 权限与知识服务"),
        ((1140, 220, 1430, 390), "DeepSeek Harness"),
        ((160, 560, 460, 740), "Semantica 解析治理"),
        ((600, 560, 900, 740), "检索 / 向量 / 图谱"),
        ((1040, 560, 1340, 740), "数据库语义查询"),
    ]
    for (x1, y1, x2, y2), label in boxes:
        draw.rounded_rectangle((x1, y1, x2, y2), radius=22, fill="#eef4ff", outline="#356ae6", width=3)
        draw.multiline_text((x1 + 25, y1 + 58), label, fill="#1f2937", font=font, align="center")
    for start, end in [((360, 305), (470, 305)), ((1030, 305), (1140, 305)), ((750, 390), (310, 560)), ((750, 390), (750, 560)), ((750, 390), (1190, 560))]:
        draw.line((*start, *end), fill="#356ae6", width=5)
    architecture = root / "NexusOne产品架构图.png"
    image.save(architecture)

    sales = EmailMessage()
    sales["Subject"] = "[验收数据] NexusOne 集团私有化方案建议"
    sales["From"] = "sales@example.test"
    sales["To"] = "customer@example.test"
    sales["Date"] = "Tue, 12 May 2026 10:00:00 +0800"
    sales.set_content(
        "NexusOne 面向集团型企业提供知识治理、混合检索、知识图谱和可追溯问答。"
        "建议先建设集团制度与产品知识，再逐步接入经营数据库和供应商场景。"
    )
    sales.add_attachment(faq.read_bytes(), maintype="text", subtype="markdown", filename=faq.name)
    sales_path = root / "NexusOne销售方案邮件.eml"
    sales_path.write_bytes(sales.as_bytes())
    return [faq, after_sales_path, architecture, sales_path]


def make_contract(root: Path) -> Path:
    from docx import Document

    doc = Document()
    add_page(doc, "关键器件采购框架协议", [
        "合同编号：GL-SC-2026-008；甲方：国联供应链有限公司；乙方：华星核心器件有限公司。",
        "合同含税金额 680000 元，履约期间为 2026年1月1日至2026年12月31日。",
        "乙方负责 NexusOne 关键器件交付，出现严重质量异常时，甲方可暂停后续采购并启动供应商复评。",
    ])
    path = root / "关键器件采购框架协议.docx"
    doc.save(path)
    return path


def make_supplier_materials(root: Path) -> list[Path]:
    from docx import Document
    from openpyxl import Workbook

    access = Document()
    add_page(access, "供应商准入申请与评审记录", [
        "申请编号：GL-SUP-ENTRY-2026-015；供应商：华星核心器件有限公司；所属区域：华东。",
        "主营产品为 NexusOne 电源与计算模块，属于关键产品部件。",
        "资质、财务和交付能力评审结果为有条件通过，要求纳入季度风险监测。",
        "准入依据：《供应商分级管理制度》V2.1 和《采购管理办法》V2.0。",
    ])
    access_path = root / "华星核心器件供应商准入材料.docx"
    access.save(access_path)

    score_book = Workbook()
    sheet = score_book.active
    sheet.title = "2026供应商评分"
    sheet.append(["供应商", "质量", "交付", "服务", "风险扣分", "综合得分", "等级"])
    sheet.append(["华星核心器件有限公司", 82, 70, 86, 20, 68, "C-重点监测"])
    sheet.append(["江南数字设备有限公司", 93, 95, 90, 0, 93, "A-优秀"])
    sheet.append(["北方云计算有限公司", 88, 90, 92, 0, 90, "A-优秀"])
    score_path = root / "2026供应商评分表.xlsx"
    score_book.save(score_path)

    minutes = Document()
    add_page(minutes, "华星核心器件风险处置会议纪要", [
        "会议编号：GL-RISK-MEETING-2026-04；会议日期：2026年4月14日。",
        "参会部门：采购管理部、风险管理部、供应商管理部、产品部。",
        "会议确认华星核心器件在 2026 年发生交付延期和质量异常两起高风险事件。",
        "决定暂停新增关键器件订单，十个工作日内完成复评；现有合同按风险条款执行。",
    ])
    minutes_path = root / "华星核心器件风险处置会议纪要.docx"
    minutes.save(minutes_path)

    catalog = Document()
    add_page(catalog, "华星核心器件产品目录", [
        "目录版本：2026-Q1；供应商：华星核心器件有限公司。",
        "NX-PWR-01 电源模块用于 NexusOne NX1，属于关键产品部件，含税单价 18000 元。",
        "NX-COMPUTE-02 计算模块用于 NexusOne NX1，属于关键产品部件，含税单价 32000 元。",
        "通用机柜附件不属于关键产品，含税单价 2800 元。",
    ])
    catalog_docx = root / "华星核心器件产品目录.docx"
    catalog.save(catalog_docx)
    catalog_pdf = convert_to_pdf(catalog_docx, root)
    catalog_docx.unlink()
    return [access_path, score_path, minutes_path, catalog_pdf]


def make_spreadsheet(root: Path) -> Path:
    from openpyxl import Workbook

    book = Workbook()
    params = book.active
    params.title = "产品参数"
    params.append(["产品编码", "产品名称", "部署方式", "标准单价", "发布年份"])
    params.append(["NX1", "NexusOne", "Docker 私有化", 100000, 2026])
    sources = book.create_sheet("支持数据源")
    sources.append(["类别", "数据源", "接入方式"])
    for category, name, method in [
        ("网页", "Web / RSS / Sitemap", "增量抓取"),
        ("研发", "Git", "版本同步"),
        ("数据库", "PostgreSQL / MySQL", "快照与实时语义查询"),
        ("对象存储", "S3 / MinIO", "对象前缀同步"),
        ("协议", "WebDAV / FTP / SFTP / IMAP / MCP", "协议连接器"),
    ]:
        sources.append([category, name, method])
    path = root / "NexusOne技术参数与数据源.xlsx"
    book.save(path)
    return path


def make_presentation(root: Path) -> Path:
    from pptx import Presentation

    deck = Presentation()
    slides = [
        ("NexusOne 产品介绍", "面向集团型企业的知识管理与智能问答一体机\n2026 企业版"),
        ("为什么需要", "统一知识沉淀\n多模态内容理解\n面向员工与智能体供给\n组织权限隔离"),
        ("核心能力", "自动治理\n混合检索\n可核验引用\n知识图谱与规则分析\nREST / MCP / CLI"),
    ]
    for title, body in slides:
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    path = root / "NexusOne产品介绍.pptx"
    deck.save(path)
    return path


def make_scanned_qualification(root: Path) -> tuple[Path, Path]:
    from PIL import Image, ImageDraw, ImageFont

    font_path = find_cjk_font()
    title_font = ImageFont.truetype(str(font_path), 48)
    body_font = ImageFont.truetype(str(font_path), 31)
    image = Image.new("RGB", (1654, 2339), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((80, 80, 1574, 2259), outline="#333333", width=4)
    draw.text((180, 180), "供应商资质审查记录", fill="black", font=title_font)
    lines = [
        "供应商：华星核心器件有限公司",
        "审查编号：GL-SUP-2026-015",
        "审查日期：2026年1月15日",
        "结论：有条件通过，纳入关键供应商持续监测。",
        "风险提示：2026年出现交付延期与质量异常两起高风险事件。",
    ]
    for index, line in enumerate(lines):
        draw.text((180, 350 + index * 110), line, fill="black", font=body_font)
    draw.text((1050, 1900), "国联供应链（验收章）", fill="#b42318", font=body_font)
    png = root / "供应商资质审查扫描件.png"
    pdf = root / "供应商资质审查扫描件.pdf"
    image.save(png)
    image.save(pdf, "PDF", resolution=150)
    return png, pdf


def make_scanned_policy_stamp(root: Path) -> Path:
    from PIL import Image, ImageDraw, ImageFont

    font_path = find_cjk_font()
    title_font = ImageFont.truetype(str(font_path), 46)
    body_font = ImageFont.truetype(str(font_path), 30)
    image = Image.new("RGB", (1654, 2339), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((90, 90, 1564, 2249), outline="#333333", width=4)
    draw.text((330, 180), "国联集团知识管理办法", fill="black", font=title_font)
    lines = [
        "发布编号：GL-KM-POLICY-001",
        "当前版本：V2.0",
        "生效日期：2026年3月1日",
        "适用范围：集团总部、所属企业及其部门",
        "制度每六个月复核一次，废止版本保留历史但不进入当前检索。",
        "本扫描件用于验证盖章制度的 OCR、版本与来源追溯。",
    ]
    for index, line in enumerate(lines):
        draw.text((170, 390 + index * 115), line, fill="black", font=body_font)
    draw.ellipse((1070, 1690, 1430, 2050), outline="#b42318", width=12)
    draw.text((1120, 1820), "验收专用章", fill="#b42318", font=body_font)
    path = root / "国联集团知识管理办法V2.0盖章扫描版.pdf"
    image.save(path, "PDF", resolution=150)
    return path


def make_email(root: Path, contract: Path) -> Path:
    message = EmailMessage()
    message["Subject"] = "[验收数据] 华星核心器件风险处置通知"
    message["From"] = "supplier-risk@example.test"
    message["To"] = "procurement@example.test"
    message["Date"] = "Mon, 13 Apr 2026 09:30:00 +0800"
    message.set_content(
        "华星核心器件在 2026 年发生交付延期和质量异常两起高风险事件。"
        "请依据采购框架协议启动供应商复评，暂停新增关键器件订单。"
    )
    message.add_attachment(
        contract.read_bytes(),
        maintype="application",
        subtype="vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=contract.name,
    )
    path = root / "供应商风险处置通知.eml"
    path.write_bytes(message.as_bytes())
    return path


def make_media(root: Path) -> list[Path]:
    sample_rate, duration = 16_000, 2
    wav_path = root / "NexusOne培训录音.wav"
    with wave.open(str(wav_path), "wb") as wav:
        wav.setparams((1, 2, sample_rate, sample_rate * duration, "NONE", "not compressed"))
        frames = bytearray()
        for index in range(sample_rate * duration):
            value = int(0.2 * 32767 * math.sin(2 * math.pi * 440 * index / sample_rate))
            frames.extend(struct.pack("<h", value))
        wav.writeframes(frames)
    mp3_path = root / "NexusOne培训录音.mp3"
    video_path = root / "NexusOne部署演示.mp4"
    run(["ffmpeg", "-y", "-loglevel", "error", "-i", str(wav_path), "-metadata", "title=NexusOne training acceptance fixture", str(mp3_path)])
    run([
        "ffmpeg", "-y", "-loglevel", "error", "-f", "lavfi", "-i", "color=c=white:s=640x360:d=3",
        "-f", "lavfi", "-i", "sine=frequency=440:duration=3", "-metadata", "title=NexusOne deployment demo",
        "-shortest", "-pix_fmt", "yuv420p", str(video_path),
    ])
    return [wav_path, mp3_path, video_path]


def generate(root: Path = ROOT) -> dict[str, object]:
    if root.exists():
        shutil.rmtree(root)
    root.mkdir(parents=True)
    manual = make_product_manual(root)
    policies = make_policy_versions(root)
    group_policies = make_group_policy_set(root)
    contract = make_contract(root)
    supplier_materials = make_supplier_materials(root)
    spreadsheet = make_spreadsheet(root)
    presentation = make_presentation(root)
    product_support = make_product_support_materials(root)
    scans = make_scanned_qualification(root)
    policy_scan = make_scanned_policy_stamp(root)
    email = make_email(root, contract)
    media = make_media(root)
    indicator = root / "集团经营指标口径.md"
    indicator.write_text(
        "# 集团经营指标口径\n\n"
        "销售额仅统计状态为 completed 的订单，取消订单不计入。\n\n"
        "同比增长率 =（本期销售额 - 上年同期销售额）/ 上年同期销售额 × 100%。\n\n"
        "目标完成率 = 已完成销售额 / 同期销售目标 × 100%。\n",
        encoding="utf-8",
    )
    risk = root / "供应商风险评估报告.md"
    risk.write_text(
        "# 供应商风险评估报告\n\n"
        "评估对象：华星核心器件。2026 年高风险事件共 2 起，分别为交付延期与质量异常。\n\n"
        "处置建议：暂停新增关键器件订单，完成复评后恢复。\n",
        encoding="utf-8",
    )
    bundle = root / "NexusOne交付资料包.zip"
    with zipfile.ZipFile(bundle, "w", zipfile.ZIP_DEFLATED) as archive:
        for path in (manual, spreadsheet, presentation, indicator):
            archive.write(path, f"NexusOne交付资料/{path.name}")
    manifest = {
        "dataset": "guolian-acceptance-v1",
        "generated_files": sorted(path.name for path in root.iterdir()),
        "authoritative_facts": {
            "sales_2026_completed": 910000,
            "sales_2025_completed": 200000,
            "yoy_percent": 355,
            "target_2026": 1400000,
            "completion_percent": 65,
            "nexusone_sales_2026": 400000,
            "high_risk_events": 2,
        },
        "synthetic": True,
    }
    (root / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return manifest


if __name__ == "__main__":
    target = Path(os.environ.get("GUOLIAN_FIXTURE_DIR", ROOT))
    print(json.dumps(generate(target), ensure_ascii=False, indent=2))
